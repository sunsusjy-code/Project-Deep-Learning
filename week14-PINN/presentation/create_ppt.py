from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_pinn_presentation():
    # 创建 PPT 对象
    prs = Presentation()

    # 定义一个辅助函数来快速添加幻灯片
    def add_slide(title, content_bullets, image_placeholder_text=None):
        slide_layout = prs.slide_layouts[1] # 使用 "Title and Content" 布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # 设置文本内容
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # 清除默认的空段落

        for bullet in content_bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            if bullet.startswith("  -"):
                p.text = bullet.replace("  -", "")
                p.level = 1
            elif bullet.startswith("    *"):
                p.text = bullet.replace("    *", "")
                p.level = 2
        
        # 如果需要放图，添加一个提示框
        if image_placeholder_text:
            left = Inches(5)
            top = Inches(2)
            width = Inches(4)
            height = Inches(3)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = f"[PLACEHOLDER: {image_placeholder_text}]"
            p.font.bold = True
            #p.font.color.rgb = None # Default black

    # =================================================================
    # Part 1: Theoretical Foundations (10 Slides)
    # =================================================================
    add_slide("PINN Study Report: Theory to Practice", 
              ["Presenter: Sun Jingyi", "Topic: PINN, L-BFGS, and Inverse Problems", "Date: 2025-11-24"])
    
    add_slide("Motivation: Why PINN?", 
              ["Solving PDEs without mesh generation.",
               "Handling data scarcity by integrating physical laws.",
               "Crucial for high-dimensional problems (e.g., Finance)."])
    
    add_slide("The PINN Architecture", 
              ["Inputs: (x, t)",
               "Outputs: u(x, t)",
               "Loss Function = MSE_Data + lambda * MSE_PDE",
               "  - Neural Network acts as a function approximator."],
              "Diagram of PINN Architecture")

    add_slide("Automatic Differentiation", 
              ["Key Enabler: torch.autograd",
               "Allows exact calculation of derivatives (u_x, u_xx) w.r.t inputs.",
               "Not numerical approximation (Finite Difference).",
               "  - Enables the 'Physics Loss' term."])

    add_slide("Optimization Challenges", 
              ["Multi-objective Optimization:",
               "  - Balancing Boundary Loss vs. PDE Residual Loss.",
               "Landscape Issue:",
               "  - 'Long-narrow valley' problem in error surface.",
               "  - Standard Gradient Descent converges slowly here."],
              "Image of Long-narrow valley from PPT")

    add_slide("First-Order Optimization (Adam)", 
              ["Formula: x_{i+1} = x_i - lr * grad(f)",
               "Pros:",
               "  - Robust to noise.",
               "  - Computationally cheap per step.",
               "Cons:",
               "  - Ignores curvature information.",
               "  - Slow convergence near the minimum."])

    add_slide("Second-Order Optimization (Newton's)", 
              ["Uses Curvature (Hessian Matrix).",
               "Formula: x_{i+1} = x_i - H^{-1} * grad(f)",
               "Pros:",
               "  - Extremely fast convergence (Quadratic).",
               "Cons:",
               "  - Calculating Inverse Hessian is expensive."])

    add_slide("Levenberg-Marquardt (LM)", 
              ["A blend of Gradient Descent and Gauss-Newton.",
               "Formula involves (H + lambda * I)^-1.",
               "Adaptive lambda:",
               "  - Large lambda -> Gradient Descent behavior.",
               "  - Small lambda -> Newton behavior."])

    add_slide("Why L-BFGS/LM for PINN?", 
              ["Parameter scale is relatively small (< 100k) compared to LLMs.",
               "Full-matrix operations (Jacobian) are feasible.",
               "Provides high precision (1e-7) required for scientific computing."])

    add_slide("Theoretical Conclusion", 
              ["First-order (Adam) is good for exploration.",
               "Second-order (L-BFGS) is superior for fine-tuning.",
               "Next: Analyzing Professor's L-BFGS Code."])

    # =================================================================
    # Part 2: Analysis of Professor's Code (10 Slides)
    # =================================================================
    add_slide("Task Overview: Elliptic Equation", 
              ["Target Equation: -Laplacian(u) = f",
               "Domain: 2D Square [-1, 1]",
               "Reference Solution: u = cos(y)",
               "Objective: Verify L-BFGS accuracy."])

    add_slide("Code Structure Analysis", 
              ["Module: Uses 'Functorch' for Jacobian calculation.",
               "Model: Fully Connected Network (2 -> 20 -> 1).",
               "Loss: Boundary MSE + Residual MSE.",
               "Optimizer: torch.optim.LBFGS."],
              "Screenshot of Code Structure")

    add_slide("L-BFGS Implementation Detail", 
              ["Key Parameters:",
               "  - history_size = 50",
               "  - line_search_fn = 'strong_wolfe'",
               "Strong Wolfe ensures sufficient decrease per step."])

    add_slide("Residual Calculation", 
              ["Using 'functorch.jacrev' and 'vmap'.",
               "Calculates exact Hessian (Laplacian) efficiently.",
               "  - f_xx + f_yy"],
              "Code Screenshot: Residual Function")

    add_slide("Sampling Strategy", 
              ["Chebyshev Points vs. Uniform Grid.",
               "Chebyshev nodes minimize interpolation error at boundaries.",
               "Essential for high-accuracy PDE solving."],
              "Scatter plot of Grid Points")

    add_slide("Experiment Results: Convergence", 
              ["Loss dropped rapidly.",
               "Final Loss: ~ 2.3e-07",
               "Converged within 1000 iterations.",
               "Validates the high-precision claim of L-BFGS."],
              "Screenshot of Terminal Output (e-07)")

    add_slide("Experiment Results: Accuracy", 
              ["L_inf Error: 0.00023 (0.02%)",
               "L_2 Error: ~ 0.0",
               "Conclusion: L-BFGS achieves machine-precision level accuracy for steady-state problems."])

    add_slide("Visualization", 
              ["Prediction vs. Exact Solution.",
               "Visually indistinguishable.",
               "Smooth gradients captured perfectly."],
              "Screenshot of 3D Surface Plot")

    add_slide("Code Critique", 
              ["Pros:",
               "  - Extremely accurate.",
               "  - Fast convergence per epoch.",
               "Cons:",
               "  - Heavy computation per step (Jacobian)."])

    add_slide("Transition to Dynamics", 
              ["L-BFGS works great for Elliptic (Steady).",
               "Challenge: What about Time-Dependent problems with Shocks?",
               "Introducing Extension I: Burgers' Equation."])

    # =================================================================
    # Part 3: Extension I - Burgers' Equation (10 Slides)
    # =================================================================
    add_slide("Extension I: Burgers' Equation", 
              ["Goal: Simulate Fluid Dynamics with Shock Waves.",
               "Equation: u_t + u*u_x = nu * u_xx",
               "Features:",
               "  - Non-linear convection (u*u_x).",
               "  - Diffusion (nu * u_xx)."])

    add_slide("Experimental Setup", 
              ["Time: t in [0, 1]",
               "Space: x in [-1, 1]",
               "Initial Condition: u = -sin(pi * x)",
               "  - Creates a collision at x=0."])

    add_slide("Optimizer Choice: Why Adam?", 
              ["Problem Landscape is Non-convex and complex due to Shocks.",
               "L-BFGS might get stuck early.",
               "Used Adam for robust exploration of the solution space."])

    add_slide("My Implementation", 
              ["Network: 2 Inputs -> Tanh -> 1 Output.",
               "Physics Loss Implementation:",
               "  - f = u_t + u*u_x - (0.01/pi)*u_xx"],
              "Code Screenshot: Physics Loss")

    add_slide("Training Dynamics", 
              ["Initial Phase: Learning the boundary.",
               "Middle Phase: Learning the slope.",
               "Final Phase: Sharpening the shock interface."],
              "Screenshot of Loss Curve")

    add_slide("Key Result: Shock Wave Formation", 
              ["Heatmap Visualization (Time vs Space).",
               "Observation:",
               "  - Smooth sine wave at t=0.",
               "  - Sharp discontinuity forms at t=0.5.",
               "  - Validates Physics Logic."],
              "Your 'Burgers_Solution.png' Heatmap")

    add_slide("Detailed Analysis", 
              ["The 'Collision':",
               "  - Top fluid moves down (Blue).",
               "  - Bottom fluid moves up (Red).",
               "  - They meet at x=0, creating a shock."])

    add_slide("Handling Gradient Pathology", 
              ["Issue: IC Loss is harder to minimize than PDE Loss.",
               "Solution: Weighted Loss Function.",
               "  - Loss = L_PDE + 20 * L_IC + 20 * L_BC"])

    add_slide("Validation", 
              ["Comparison with Exact Solution (if available) or physical intuition.",
               "The sharpness of the shock depends on viscosity (nu)."])

    add_slide("Section Summary", 
              ["Successfully extended PINN to time-dependent dynamics.",
               "Simulated Shock Waves using Adam.",
               "Next: Can we inverse this process?"])

    # =================================================================
    # Part 4: Extension II - Inverse Problem (10 Slides)
    # =================================================================
    add_slide("Extension II: Inverse Problem", 
              ["Motivation: Quantitative Finance (Model Calibration).",
               "Scenario: We see prices (u), but don't know Volatility (nu).",
               "Goal: Discover 'nu' from noisy data."])

    add_slide("Mathematical Formulation", 
              ["Assumption: u_t + u*u_x = [?] * u_xx",
               "Treat 'nu' as a learnable parameter.",
               "Loss = L_Data + L_PDE(nu)"])

    add_slide("Implementation Strategy", 
              ["Code: self.nu_log = nn.Parameter(...)",
               "Using Logarithm to ensure positivity.",
               "Optimizer updates both Weights and Nu."],
              "Code Screenshot: nn.Parameter")

    add_slide("Case Study 1: The Failure", 
              ["Initial Attempt: No Observation Data.",
               "Result: Nu converged to 0.15 (Wrong!) instead of 0.003.",
               "Reason: 'Spectral Bias'. Network prefers smooth solutions (High Nu)."],
              "Screenshot of Failed Nu Plot")

    add_slide("The Importance of Data Quality", 
              ["Garbage In, Garbage Out.",
               "Generated high-quality 'Shock Data' using a pre-trained Teacher model.",
               "Constraint: Data forces the network to accept the sharp shock."])

    add_slide("Optimization Innovation: Parameter Grouping", 
              ["Challenge: Nu converges too slowly compared to weights.",
               "Solution: Differential Learning Rates.",
               "  - Net LR = 0.005",
               "  - Nu LR = 0.05 (10x faster)"],
              "Code Screenshot: Optimizer Grouping")

    add_slide("Success Result: Convergence", 
              ["Observed Nu dropping from 0.36 to 0.003.",
               "Monotonic convergence.",
               "Demonstrates robust calibration capability."],
              "Your 'Nu_Fixed.png' Plot")

    add_slide("Accuracy Assessment", 
              ["Initial Error: > 1000%",
               "Final Error: < 15%",
               "Achieved correct order of magnitude (1e-3)."])

    add_slide("Financial Analogy", 
              ["Physics: Viscosity (nu) -> Smoothness of Shock.",
               "Finance: Volatility (sigma) -> Spread of Probability.",
               "This experiment simulates 'Implied Volatility' calibration."])

    add_slide("Section Summary", 
              ["Successfully solved the Inverse Problem.",
               "Highlighted the need for good data and specific optimization strategies."])

    # =================================================================
    # Part 5: Discussion & Conclusion (8 Slides)
    # =================================================================
    add_slide("Deep Comparison: Adam vs. L-BFGS", 
              ["Precision: L-BFGS (1e-7) >>> Adam (1e-3).",
               "Robustness: Adam >>> L-BFGS (better for noise).",
               "Speed: L-BFGS is faster per epoch, but heavier per step."])

    add_slide("Trade-off Analysis", 
              ["When to use L-BFGS?",
               "  - Forward problems, Smooth solutions, High precision needed.",
               "When to use Adam?",
               "  - Inverse problems, Noisy data, Complex landscapes (Shocks)."])

    add_slide("Proposed Hybrid Strategy", 
              ["Phase 1: Adam",
               "  - Fast initialization, finding the 'basin' of global minimum.",
               "Phase 2: L-BFGS",
               "  - Fine-tuning to reach machine precision (1e-7)."])

    add_slide("Application in Quant Finance: Pricing", 
              ["Solving Black-Scholes PDE.",
               "Handling Exotic Options (Barrier, Asian) where formulas don't exist.",
               "High-dimensional Basket Options."])

    add_slide("Application in Quant Finance: Calibration", 
              ["Constructing Local Volatility Surfaces.",
               "Using Inverse PINN to match market option prices.",
               "Real-time parameter estimation."])

    add_slide("Future Work", 
              ["High-Dimensionality:",
               "  - Curse of dimensionality in PDEs.",
               "  - Solution: Monte-Carlo sampling + PINN.",
               "Transfer Learning:",
               "  - Re-using models for changing market conditions."])

    add_slide("Conclusion", 
              ["Verified Professor's L-BFGS code (High Precision).",
               "Extended to Time-Dependent Shock Waves (Project A).",
               "Solved Inverse Calibration Problem (Project B).",
               "Gained deep insight into Optimization Strategies."])

    add_slide("Q & A", ["Thank you for listening."])

    # 保存文件
    prs.save('PINN_Final_Report.pptx')
    print("✅ 成功生成 PPT! 文件名: PINN_Final_Report.pptx")
    print("👉 请现在去文件夹里打开它，把你之前跑出来的截图拖进去填空！")

if __name__ == "__main__":
    create_pinn_presentation()